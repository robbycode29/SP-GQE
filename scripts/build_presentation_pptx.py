#!/usr/bin/env python3
"""Build SP-GQE dissertation presentation (motivation, method, results, visuals)."""

from __future__ import annotations

import json
import sys
from pathlib import Path

REPO = Path(__file__).resolve().parents[1]
SCRIPTS = Path(__file__).resolve().parent
sys.path.insert(0, str(SCRIPTS))

from build_research_reports_docx import (  # noqa: E402
    QUERIES_SAMPLE,
    RES,
    _extract_sparql_samples,
    _load_summary,
)

OUT_DIR = REPO / "deliverables" / "rapoarte-cercetare"
ASSETS = OUT_DIR / "_ppt_assets"
OUT_PPTX = OUT_DIR / "Robert Oprescu - SP-GQE Presentation.pptx"

# Slide palette
C_NAVY = (26, 54, 93)
C_BLUE = (43, 108, 176)
C_ORANGE = (221, 107, 32)
C_LIGHT = (237, 242, 247)
C_WHITE = (255, 255, 255)
C_GRAY = (113, 128, 150)


def _rgb(t: tuple[int, int, int]):
    from pptx.dml.color import RGBColor

    return RGBColor(*t)


def _ensure_assets() -> None:
    ASSETS.mkdir(parents=True, exist_ok=True)


def create_pipeline_flowchart(path: Path) -> Path:
    """Top-down branching logic diagram (matplotlib)."""
    import matplotlib.pyplot as plt
    from matplotlib.patches import FancyBboxPatch

    fig, ax = plt.subplots(figsize=(10, 7), dpi=150)
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 10)
    ax.axis("off")

    def box(x, y, w, h, text, fc="#ebf8ff", ec="#2b6cb0", fs=9):
        p = FancyBboxPatch(
            (x, y),
            w,
            h,
            boxstyle="round,pad=0.02,rounding_size=0.08",
            facecolor=fc,
            edgecolor=ec,
            linewidth=1.5,
        )
        ax.add_patch(p)
        ax.text(
            x + w / 2,
            y + h / 2,
            text,
            ha="center",
            va="center",
            fontsize=fs,
            wrap=True,
        )

    def arrow(x1, y1, x2, y2):
        ax.annotate(
            "",
            xy=(x2, y2),
            xytext=(x1, y1),
            arrowprops=dict(arrowstyle="->", color="#4a5568", lw=1.5),
        )

    box(3.2, 8.6, 3.6, 0.9, "HotpotQA question q\n(10 distractor paragraphs)", fc="#bee3f8")
    box(3.0, 7.2, 4.0, 0.8, "Build RDF graph G(q)\nEntity + coOccurs triples", fc="#e6fffa")

    box(0.4, 5.5, 2.6, 1.0, "Branch 1\nSPARQL n-hop\nfrom NER seeds → A", fc="#fefcbf")
    box(3.7, 5.5, 2.6, 1.0, "Branch 2\nSPARQL keyword\non rdfs:label → B", fc="#fefcbf")
    box(7.0, 5.5, 2.6, 1.0, "Seeds + probes\n(spaCy NER / chunks)", fc="#edf2f7", fs=8)

    box(2.5, 3.8, 5.0, 0.9, "Fusion: candidates = (A ∪ B) \\ seeds", fc="#faf089")
    box(2.2, 2.5, 5.6, 0.9, "Semantic prune: max cosine vs {q}∪probes\nkeep if score ≥ τ (MiniLM)", fc="#faf089")

    box(2.0, 1.1, 2.8, 0.9, "Augmented FAISS\nquery → top-k chunks", fc="#c6f6d5")
    box(5.2, 1.1, 2.8, 0.9, "Groq reader\nllama-3.1-8b (T=0)", fc="#c6f6d5")

    arrow(5, 8.6, 5, 8.0)
    arrow(5, 7.2, 1.7, 6.5)
    arrow(5, 7.2, 5, 6.5)
    arrow(5, 7.2, 8.3, 6.5)
    arrow(1.7, 5.5, 3.5, 4.7)
    arrow(5, 5.5, 5, 4.7)
    arrow(8.3, 5.5, 6.5, 4.7)
    arrow(5, 3.8, 5, 3.4)
    arrow(3.4, 2.5, 3.4, 2.0)
    arrow(6.6, 2.5, 6.6, 2.0)

    ax.text(
        5,
        0.35,
        "V-RAG baseline skips graph branches (FAISS on question only)",
        ha="center",
        fontsize=8,
        color="#718096",
        style="italic",
    )
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return path


def create_validity_chart(path: Path, gv: dict) -> Path:
    import matplotlib.pyplot as plt

    stages = ["Branch 1\n(n-hop)", "Branch 2\n(keyword)", "Union", "Kept @τ=0.5"]
    prec = [
        gv["branch1_nhop"]["mean_precision"],
        gv["branch2_keyword"]["mean_precision"],
        gv["union"]["mean_precision"],
        gv["kept_after_tau"]["mean_precision"],
    ]
    rec = [
        gv["branch1_nhop"]["mean_recall"],
        gv["branch2_keyword"]["mean_recall"],
        gv["union"]["mean_recall"],
        gv["kept_after_tau"]["mean_recall"],
    ]
    x = range(len(stages))
    w = 0.35
    fig, ax = plt.subplots(figsize=(8, 4.5), dpi=150)
    ax.bar([i - w / 2 for i in x], prec, w, label="Precision", color="#2b6cb0")
    ax.bar([i + w / 2 for i in x], rec, w, label="Recall", color="#dd6b20")
    ax.set_xticks(list(x))
    ax.set_xticklabels(stages)
    ax.set_ylim(0, 0.75)
    ax.set_ylabel("Mean P / R")
    ax.set_title(f"Graph-query validity (n={gv['n_questions']} questions)")
    ax.legend()
    ax.grid(axis="y", alpha=0.3)
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return path


def create_pipeline_f1_chart(path: Path, agg: dict) -> Path:
    import matplotlib.pyplot as plt

    order = [
        "GR-RAG",
        "V-RAG",
        "GQE-RAG(n=2)",
        "SP-GQE(n=2,τ=0.5)",
        "GF-RAG",
        "SP-GQE-i(n=3,τ=0.5)",
    ]
    labels = [
        "GR-RAG",
        "V-RAG",
        "GQE-RAG",
        "SP-GQE\n(2,0.5)",
        "GF-RAG",
        "SP-GQE-i",
    ]
    means = [agg[k]["mean_f1"] for k in order]
    colors = ["#2b6cb0" if "SP-GQE" not in k else "#dd6b20" for k in order]
    fig, ax = plt.subplots(figsize=(8, 4.5), dpi=150)
    bars = ax.bar(labels, means, color=colors, edgecolor="#1a365d", linewidth=0.8)
    ax.set_ylabel("Mean token F1")
    ax.set_ylim(0.5, 0.6)
    ax.set_title("Per-pipeline mean F1 (28 seeds, 680 questions)")
    for b, m in zip(bars, means):
        ax.text(b.get_x() + b.get_width() / 2, m + 0.002, f"{m:.3f}", ha="center", fontsize=9)
    ax.grid(axis="y", alpha=0.3)
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return path


def _blank_slide(prs, title: str):
    from pptx.util import Inches, Pt

    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.55))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(C_NAVY)
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.08), Inches(9), Inches(0.45))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = _rgb(C_WHITE)
    return slide


def _add_stat_box(slide, left, top, width, height, value: str, label: str):
    from pptx.util import Inches, Pt

    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(C_LIGHT)
    shape.line.color.rgb = _rgb(C_BLUE)
    tf = shape.text_frame
    tf.clear()
    p0 = tf.paragraphs[0]
    p0.text = value
    p0.font.size = Pt(28)
    p0.font.bold = True
    p0.font.color.rgb = _rgb(C_ORANGE)
    p0.alignment = 1
    p1 = tf.add_paragraph()
    p1.text = label
    p1.font.size = Pt(11)
    p1.font.color.rgb = _rgb(C_NAVY)
    p1.alignment = 1


def _add_bullets(slide, top_in: float, bullets: list[str], width=9.0):
    from pptx.util import Inches, Pt

    tb = slide.shapes.add_textbox(Inches(0.55), Inches(top_in), Inches(width), Inches(5.5))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(16)
        p.space_after = Pt(8)


def _add_image(slide, path: Path, left, top, width, caption: str | None = None):
    from pptx.util import Inches, Pt

    if path.is_file():
        slide.shapes.add_picture(str(path), left, top, width=width)
    if caption:
        tb = slide.shapes.add_textbox(left, top + width * 0.55, Inches(9), Inches(0.5))
        p = tb.text_frame.paragraphs[0]
        p.text = caption
        p.font.size = Pt(10)
        p.font.italic = True
        p.font.color.rgb = _rgb(C_GRAY)


def _add_table_slide(prs, title: str, headers: list[str], rows: list[tuple]):
    from pptx.util import Inches, Pt

    slide = _blank_slide(prs, title)
    nrows, ncols = len(rows) + 1, len(headers)
    tbl = slide.shapes.add_table(nrows, ncols, Inches(0.5), Inches(1.0), Inches(9), Inches(0.4 * nrows)).table
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(C_NAVY)
        for p in cell.text_frame.paragraphs:
            p.font.color.rgb = _rgb(C_WHITE)
            p.font.size = Pt(11)
            p.font.bold = True
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
    return slide


def _sparql_short(sparql: str, max_lines: int = 7) -> str:
    lines = sparql.strip().splitlines()
    if len(lines) <= max_lines:
        return sparql.strip()
    return "\n".join(lines[:max_lines]) + "\n…"


def _add_sparql_slide(prs, samp: dict, idx: int) -> None:
    from pptx.util import Inches, Pt

    title = f"SPARQL sample {idx}: [{samp['type']}]"
    slide = _blank_slide(prs, title)

    q = samp["question"]
    if len(q) > 100:
        q = q[:97] + "…"
    meta = (
        f"Gold: {samp.get('gold_answer', '—')}  |  "
        f"F1 SP-GQE / V-RAG: {samp.get('f1_sp_gqe', '—')} / {samp.get('f1_v_rag', '—')}"
    )
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.65), Inches(9), Inches(0.9))
    tf = tb.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = q
    p0.font.size = Pt(12)
    p0.font.bold = True
    p1 = tf.add_paragraph()
    p1.text = meta
    p1.font.size = Pt(10)
    p1.font.color.rgb = _rgb(C_GRAY)

    y = 1.35
    for bi, br in enumerate(samp.get("branches", [])[:2]):
        tb = slide.shapes.add_textbox(Inches(0.45), Inches(y), Inches(4.5), Inches(2.2))
        tf = tb.text_frame
        tf.word_wrap = True
        h = tf.paragraphs[0]
        h.text = f"Branch {bi + 1}: {br['label']}"
        h.font.size = Pt(10)
        h.font.bold = True
        code = tf.add_paragraph()
        code.text = _sparql_short(br.get("sparql", ""))
        code.font.name = "Consolas"
        code.font.size = Pt(7)
        res = tf.add_paragraph()
        res.text = f"Results (n={br.get('returned_n', '?')}): {br.get('returned_entities', '')[:200]}"
        res.font.size = Pt(8)
        y += 2.35

    # Fusion column
    fusion_lines = []
    if samp.get("union_size"):
        fusion_lines.append(f"Union |A∪B| = {samp['union_size']}")
    if samp.get("kept_n"):
        fusion_lines.append(f"Kept @τ=0.5 (n={samp['kept_n']}): {samp.get('kept_entities', '')[:180]}")
    if samp.get("top_similarities"):
        top = ", ".join(f"{e} ({s})" for e, s in samp["top_similarities"][:4])
        fusion_lines.append(f"Top cosine: {top}")
    if samp.get("validity_rows"):
        v = "  |  ".join(f"{s}: P={p} R={r}" for s, p, r in samp["validity_rows"])
        fusion_lines.append(f"Validity: {v}")

    tb2 = slide.shapes.add_textbox(Inches(5.0), Inches(1.35), Inches(4.5), Inches(5.5))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "Fusion, pruning & validity"
    p.font.bold = True
    p.font.size = Pt(11)
    for line in fusion_lines:
        bp = tf2.add_paragraph()
        bp.text = line
        bp.font.size = Pt(9)


def build_presentation() -> Path:
    from pptx import Presentation
    from pptx.util import Inches

    summary = _load_summary()
    agg = summary["aggregated_across_seeds"]
    merged = summary["merged_config_aggregates"]
    sel = merged["selection"]
    sp_sel = merged["sp_gqe_at_selected_config"]
    proto = summary["protocol_default_n2_tau05"]
    gv = proto["graph_query_validity_pooled"]
    n_q = summary["total_question_instances"]
    n_seeds = summary["n_files"]

    _ensure_assets()
    flow_png = create_pipeline_flowchart(ASSETS / "pipeline_flow.png")
    validity_png = create_validity_chart(ASSETS / "validity_bars.png", gv)
    f1_png = create_pipeline_f1_chart(ASSETS / "pipeline_f1_bars.png", agg)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "SP-GQE: Semantic Pruning on Graph-structured Query Expansion"
    slide.placeholders[1].text = (
        "Robert Oprescu\n"
        "Motivation · Method · Results\n"
        f"HotpotQA distractor — {n_q} questions · {n_seeds} seeds"
    )

    # Outline
    s = _blank_slide(prs, "Outline")
    _add_bullets(
        s,
        0.75,
        [
            "Motivation & hypothesis (H1)",
            "Method: per-question RDF graph + two SPARQL branches",
            "Top-down pipeline logic (fusion & semantic pruning)",
            "Fair baselines & evaluation design",
            f"Results at scale ({n_q} questions, {n_seeds} seeds)",
            "Heatmap (n, τ) selection & graph-query validity",
            "SPARQL query + result samples",
        ],
    )

    # Motivation
    s = _blank_slide(prs, "Motivation")
    _add_bullets(
        s,
        0.7,
        [
            "Multi-hop QA needs evidence across paragraphs — dense RAG may miss bridge facts.",
            "Graph-augmented systems vary in where structure acts (index vs query time).",
            "SP-GQE: lightweight, training-free, query-time SPARQL on an in-corpus co-occurrence graph.",
            "No external KB — only HotpotQA’s 10 distractor paragraphs per question.",
            "Fair comparison: same FAISS index, MiniLM embeddings, Groq reader for all pipelines.",
        ],
    )
    _add_stat_box(s, Inches(6.8), Inches(1.2), Inches(2.8), Inches(1.1), str(n_q), "questions evaluated")
    _add_stat_box(s, Inches(6.8), Inches(2.5), Inches(2.8), Inches(1.1), str(n_seeds), "random seeds")
    _add_stat_box(s, Inches(6.8), Inches(3.8), Inches(2.8), Inches(1.1), "50%", "bridge / comparison stratified")

    # Hypothesis
    s = _blank_slide(prs, "Hypothesis H1")
    _add_bullets(
        s,
        0.7,
        [
            "H1: SP-GQE improves token F1 vs V-RAG on bridge questions (multi-hop linking).",
            "Null: on comparison questions, dense retrieval often suffices.",
            "Diagnostics: Do branches complement? Does τ pruning trade P for R?",
        ],
    )

    # Pipeline flowchart
    s = _blank_slide(prs, "Method — top-down pipeline (two-branch logic)")
    _add_image(s, flow_png, Inches(0.35), Inches(0.7), Inches(9.3))

    # Graph + branches text
    s = _blank_slide(prs, "Method — RDF graph & SPARQL branches")
    _add_bullets(
        s,
        0.7,
        [
            "Graph: spaCy entities + symmetric spg:coOccurs within sentences + rdfs:label.",
            "Branch 1: bounded property paths from NER seeds (n hops) → set A.",
            "Branch 2: CONTAINS filters on labels from noun-chunk keywords → set B.",
            "Default protocol: n=2, τ=0.5; heatmap sweep n∈{1,2,3}, τ∈{0.3…0.7}.",
        ],
    )

    # Baselines
    _add_table_slide(
        prs,
        "Baselines (same index + reader)",
        ["Pipeline", "Graph", "Role"],
        [
            ("V-RAG", "—", "Dense retrieval on question only"),
            ("GQE-RAG (n=2)", "Branch 1", "Structural expansion, no τ prune"),
            ("SP-GQE (n=2, τ=0.5)", "A ∪ B + prune", "Primary system"),
            ("SP-GQE-i (n=3, τ=0.5)", "Iterative B1", "Ablated variant"),
            ("GR-RAG / GF-RAG", "Lexical", "Graph-adjacent controls"),
        ],
    )

    # Evaluation
    s = _blank_slide(prs, "Evaluation design")
    _add_bullets(
        s,
        0.7,
        [
            "Primary metric: token F1; secondary: EM, supporting-title recall@k, P@k.",
            "Paired ΔF1: SP-GQE − V-RAG with bootstrap 95% CI (bridge / comparison).",
            "Graph-query validity: P/R vs gold supporting NER at 4 stages.",
            "Aggregation: mean ± 95% CI across seed-level means; multi-day daily_runs JSON.",
        ],
    )
    _add_stat_box(s, Inches(6.5), Inches(1.5), Inches(3.0), Inches(1.0), "680", "question instances (full study)")
    _add_stat_box(s, Inches(6.5), Inches(2.7), Inches(3.0), Inches(1.0), "562", "questions (semester preliminary)")

    # Results headline
    s = _blank_slide(prs, "Results — experiment scale")
    _add_stat_box(s, Inches(0.6), Inches(1.2), Inches(2.2), Inches(1.2), str(n_q), "questions")
    _add_stat_box(s, Inches(3.0), Inches(1.2), Inches(2.2), Inches(1.2), str(n_seeds), "seeds")
    _add_stat_box(s, Inches(5.4), Inches(1.2), Inches(2.2), Inches(1.2), "24", "heatmap seeds")
    _add_stat_box(s, Inches(7.8), Inches(1.2), Inches(2.0), Inches(1.2), "580", "heatmap weight")
    _add_bullets(
        s,
        2.8,
        [
            "Reader: Groq llama-3.1-8b-instant (T=0). Embeddings: all-MiniLM-L6-v2.",
            f"Merged-grid best: n={sel['n_hops']}, τ={sel['tau']} "
            f"(pooled F1={sel['pooled_mean_f1_on_merged_grid']:.3f}).",
            "Protocol-default SP-GQE (2, 0.5) used for validity & appendix paired tests.",
        ],
    )

    # F1 chart
    s = _blank_slide(prs, "Results — mean token F1 by pipeline")
    _add_image(
        s,
        f1_png,
        Inches(0.8),
        Inches(0.85),
        Inches(8.5),
        f"28 seeds, {n_q} questions; error bars = 95% CI on seed means in aggregate tables",
    )

    # Pipeline table
    _add_table_slide(
        prs,
        "Results — pipeline table (28 seeds)",
        ["Pipeline", "Mean F1", "95% CI", "P@k"],
        [
            ("GR-RAG", f"{agg['GR-RAG']['mean_f1']:.4f}", "[0.536, 0.600]", f"{agg['GR-RAG']['mean_retrieval_p_at_k']:.4f}"),
            ("V-RAG", f"{agg['V-RAG']['mean_f1']:.4f}", "[0.531, 0.594]", f"{agg['V-RAG']['mean_retrieval_p_at_k']:.4f}"),
            ("GQE-RAG(n=2)", f"{agg['GQE-RAG(n=2)']['mean_f1']:.4f}", "[0.527, 0.597]", f"{agg['GQE-RAG(n=2)']['mean_retrieval_p_at_k']:.4f}"),
            ("SP-GQE(n=2,τ=0.5)", f"{agg['SP-GQE(n=2,τ=0.5)']['mean_f1']:.4f}", "[0.528, 0.584]", f"{agg['SP-GQE(n=2,τ=0.5)']['mean_retrieval_p_at_k']:.4f}"),
            (f"SP-GQE(n={sel['n_hops']},τ={sel['tau']})*", f"{sp_sel['mean_f1']:.4f}", "[0.550, 0.605]", f"{sp_sel['mean_retrieval_p_at_k']:.4f}"),
        ],
    )

    # Paired delta
    pd = proto["paired_delta_f1_SP_GQE_minus_V_RAG"]
    _add_table_slide(
        prs,
        "Results — SP-GQE(2,0.5) − V-RAG paired ΔF1",
        ["Subset", "Mean Δ", "Bootstrap 95% CI", "n pairs"],
        [
            ("bridge", f"{pd['bridge']['mean']:.4f}", "[-0.028, 0.021]", str(pd["bridge"]["n_pairs"])),
            ("comparison", f"{pd['comparison']['mean']:.4f}", "[-0.038, 0.020]", str(pd["comparison"]["n_pairs"])),
        ],
    )
    from pptx.util import Pt

    s = prs.slides[-1]
    tb = s.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(0.6))
    tb.text_frame.paragraphs[0].text = (
        "CIs include zero → no statistically reliable F1 gain vs V-RAG at this scale (H1 not confirmed)."
    )
    tb.text_frame.paragraphs[0].font.size = Pt(12)

    # Heatmap images
    merged_f1 = RES / "merged_heatmap" / "heatmap_merged_mean_f1.png"
    merged_pk = RES / "merged_heatmap" / "heatmap_merged_mean_p_at_k.png"
    s = _blank_slide(prs, f"Results — merged (n, τ) heatmaps (weight={sel['merged_total_weight']})")
    if merged_f1.is_file():
        s.shapes.add_picture(str(merged_f1), Inches(0.4), Inches(0.75), width=Inches(4.6))
    if merged_pk.is_file():
        s.shapes.add_picture(str(merged_pk), Inches(5.0), Inches(0.75), width=Inches(4.6))
    cap = s.shapes.add_textbox(Inches(0.4), Inches(6.5), Inches(9.2), Inches(0.5))
    cap.text_frame.paragraphs[0].text = (
        f"Argmax pooled F1 → n*={sel['n_hops']}, τ*={sel['tau']} "
        f"(SP-GQE mean F1={sp_sel['mean_f1']:.4f} on {sp_sel['n_seeds']} heatmap seeds)"
    )

    # Validity
    s = _blank_slide(prs, f"Results — graph-query validity ({gv['n_questions']} questions)")
    _add_image(s, validity_png, Inches(1.0), Inches(0.9), Inches(8.0))
    _add_bullets(
        s,
        5.5,
        [
            "Union recall > Branch 2 alone → branches are complementary.",
            "τ=0.5 increases precision, lowers recall → pruning trade-off.",
        ],
    )

    # SPARQL samples
    if QUERIES_SAMPLE.is_file():
        samples = _extract_sparql_samples(QUERIES_SAMPLE, max_questions=2)
        for i, samp in enumerate(samples, start=1):
            _add_sparql_slide(prs, samp, i)

    # Conclusions
    s = _blank_slide(prs, "Conclusions")
    _add_bullets(
        s,
        0.7,
        [
            f"Rigorous multi-seed study ({n_q} Q, {n_seeds} seeds) under fair pipeline controls.",
            "SP-GQE mechanism is inspectable (SPARQL + validity stages); branches complement at entity level.",
            "Mean F1: GR-RAG / GQE-RAG competitive; protocol SP-GQE does not beat V-RAG reliably.",
            f"Tuned grid (n*={sel['n_hops']}, τ*={sel['tau']}) raises heatmap-seed F1 but paired CIs still span zero.",
            "Future: per-question logging at (n*, τ*), SP-GQE-i heatmap arm, larger samples.",
        ],
    )

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    try:
        prs.save(str(OUT_PPTX))
        return OUT_PPTX
    except PermissionError:
        alt = OUT_PPTX.with_stem(OUT_PPTX.stem + " (regenerated)")
        prs.save(str(alt))
        print(f"WARNING: could not overwrite {OUT_PPTX}; wrote {alt}")
        return alt


def main() -> None:
    try:
        import pptx  # noqa: F401
    except ImportError:
        import subprocess

        subprocess.check_call(
            [sys.executable, "-m", "pip", "install", "python-pptx>=1.0.0", "matplotlib>=3.8.0"]
        )
    path = build_presentation()
    print(f"Wrote {path}")


if __name__ == "__main__":
    main()
